# fileio.py — presentation import/export adapters (in-process Python libraries).
# SPDX-License-Identifier: GPL-3.0-or-later
#
# Maps a deck (list of Fabric.js slide JSON) <-> .pptx (python-pptx) / .odp (odfpy).
# Text boxes and pictures round-trip; positions scale between EMU and the 960x540 canvas.

import base64
import os
from io import BytesIO

CANVAS_W = 960
CANVAS_H = 540


def _ext(path):
    return os.path.splitext(path)[1].lower()


def read_deck(path):
    ext = _ext(path)
    if ext == '.pptx':
        return _read_pptx(path)
    if ext == '.odp':
        return _read_odp(path)
    raise ValueError(f'Unsupported format: {ext}')


def write_deck(path, slides):
    ext = _ext(path)
    if ext == '.pptx':
        return _write_pptx(path, slides)
    if ext == '.odp':
        return _write_odp(path, slides)
    raise ValueError(f'Unsupported format: {ext}')


def _slide(objects):
    return {'version': '5.5.2', 'objects': objects, 'background': '#ffffff'}


# ----- pptx -----------------------------------------------------------------

def _read_pptx(path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    slides = []
    for slide in prs.slides:
        objects = []
        for shape in slide.shapes:
            left = (shape.left or 0) / sw * CANVAS_W
            top = (shape.top or 0) / sh * CANVAS_H
            width = (shape.width or 0) / sw * CANVAS_W
            height = (shape.height or 0) / sh * CANVAS_H
            if shape.has_text_frame and shape.text_frame.text:
                objects.append({
                    'type': 'i-text', 'text': shape.text_frame.text,
                    'left': left, 'top': top, 'width': max(width, 50),
                    'fontSize': 32, 'fontFamily': 'sans-serif', 'fill': '#202020',
                })
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                src = 'data:%s;base64,%s' % (
                    img.content_type, base64.b64encode(img.blob).decode())
                objects.append({
                    'type': 'image', 'src': src, 'left': left, 'top': top,
                    'width': width, 'height': height, 'scaleX': 1, 'scaleY': 1,
                })
        slides.append(_slide(objects))
    return slides or [_slide([])]


def _write_pptx(path, slides):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    prs = Presentation()
    sw, sh = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    for slide in slides or [_slide([])]:
        s = prs.slides.add_slide(blank)
        for obj in (slide or {}).get('objects', []):
            left = Emu(int(obj.get('left', 0) / CANVAS_W * sw))
            top = Emu(int(obj.get('top', 0) / CANVAS_H * sh))
            width = Emu(int(obj.get('width', 100) * obj.get('scaleX', 1) / CANVAS_W * sw))
            height = Emu(int(obj.get('height', 40) * obj.get('scaleY', 1) / CANVAS_H * sh))
            kind = obj.get('type')
            if kind in ('i-text', 'text', 'textbox'):
                box = s.shapes.add_textbox(left, top, width or Emu(914400), height or Emu(457200))
                tf = box.text_frame
                tf.word_wrap = True
                tf.text = obj.get('text', '')
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].font.size = Pt(int(obj.get('fontSize', 24) * 0.75))
            elif kind == 'image':
                src = obj.get('src', '')
                if src.startswith('data:') and ',' in src:
                    blob = base64.b64decode(src.split(',', 1)[1])
                    s.shapes.add_picture(BytesIO(blob), left, top, width or None, height or None)
    prs.save(path)


# ----- odp ------------------------------------------------------------------

def _read_odp(path):
    from odf.opendocument import load
    from odf.draw import Page, Frame
    from odf.text import P
    doc = load(path)
    slides = []
    for page in doc.getElementsByType(Page):
        objects = []
        for frame in page.getElementsByType(Frame):
            text = ''.join(str(p) for p in frame.getElementsByType(P))
            if text:
                objects.append({
                    'type': 'i-text', 'text': text, 'left': 80, 'top': 80 + 60 * len(objects),
                    'width': 600, 'fontSize': 32, 'fontFamily': 'sans-serif', 'fill': '#202020',
                })
        slides.append(_slide(objects))
    return slides or [_slide([])]


def _write_odp(path, slides):
    from odf.opendocument import OpenDocumentPresentation
    from odf.draw import Page, Frame, TextBox
    from odf.text import P
    doc = OpenDocumentPresentation()
    for slide in slides or [_slide([])]:
        page = Page()
        for obj in (slide or {}).get('objects', []):
            if obj.get('type') in ('i-text', 'text', 'textbox'):
                x = '%dpx' % int(obj.get('left', 0))
                y = '%dpx' % int(obj.get('top', 0))
                w = '%dpx' % int(obj.get('width', 300) * obj.get('scaleX', 1))
                frame = Frame(x=x, y=y, width=w, height='80px')
                tb = TextBox()
                tb.addElement(P(text=obj.get('text', '')))
                frame.addElement(tb)
                page.addElement(frame)
        doc.presentation.addElement(page)
    doc.save(path)
